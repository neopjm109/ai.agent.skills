#!/usr/bin/env python3
"""Extract slide text and speaker notes from a .pptx deck.

Prints each slide as a `=== SLIDE n ===` block followed by its text lines, and
speaker notes as `[NOTES] ...`. Prefers python-pptx; falls back to reading the
raw slide XML (a .pptx is a ZIP) using only the standard library, so the sole
requirement is python3.

Usage: python3 scripts/extract.py <file.pptx>
"""
import sys


def with_python_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, 1):
        print(f"=== SLIDE {i} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        print(text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                print(f"[NOTES] {notes}")


def with_stdlib(path):
    import re
    import zipfile
    from xml.etree import ElementTree as ET

    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    def texts(xml):
        root = ET.fromstring(xml)
        return [t.text for t in root.iter(f"{A}t") if t.text and t.text.strip()]

    with zipfile.ZipFile(path) as z:
        slides = sorted(
            (n for n in z.namelist()
             if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)),
            key=lambda n: int(re.search(r"(\d+)", n).group(1)),
        )
        for i, name in enumerate(slides, 1):
            print(f"=== SLIDE {i} ===")
            for line in texts(z.read(name)):
                print(line)
            note = name.replace("slides/slide", "notesSlides/notesSlide")
            if note in z.namelist():
                joined = " ".join(texts(z.read(note))).strip()
                if joined:
                    print(f"[NOTES] {joined}")


def main():
    if len(sys.argv) != 2:
        sys.exit("usage: extract.py <file.pptx>")
    path = sys.argv[1]
    try:
        with_python_pptx(path)
    except ImportError:
        with_stdlib(path)


if __name__ == "__main__":
    main()
